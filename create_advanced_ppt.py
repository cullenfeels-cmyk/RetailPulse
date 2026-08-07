from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_advanced_retailpulse_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- UNIQUE & PROFESSIONAL COLOR PALETTE (Cyber-Executive Theme) ---
    COLOR_BG_DARK = RGBColor(15, 23, 42)      # Deep Midnight Slate
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # Ultra Light Ice Gray
    COLOR_CARD_BG = RGBColor(255, 255, 255)   # Crisp White Card Background
    COLOR_ACCENT_BLUE = RGBColor(37, 99, 235) # Electric Tech Blue
    COLOR_ACCENT_CYAL = RGBColor(14, 165, 233)# Bright Cyan Highlight
    COLOR_TEXT_DARK = RGBColor(30, 41, 59)    # Deep Slate Text
    COLOR_TEXT_LIGHT = RGBColor(255, 255, 255)# White Text
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Muted Gray Text
    COLOR_CARD_BORDER = RGBColor(226, 232, 240) # Subtle Card Border

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_advanced_header(slide, title_text, category_text="RETAILPULSE ENTERPRISE AI"):
        # Top category tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT_BLUE

        # Main Action Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_DARK

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # ==================== SLIDE 1: Dark Executive Title Slide ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_BG_DARK)

    # Decorative tech accent line
    accent_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(1.5), Inches(0.08))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_ACCENT_CYAL
    accent_line.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.5), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "RetailPulse"
    p.font.name = 'Arial'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT

    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Powered Customer Analytics, Demand Forecasting & MLOps Intelligence Platform"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT_CYAL
    p_sub.space_before = Pt(12)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Zidio Internship Executive Technical Presentation | Advanced Production Grade Architecture"
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_desc.space_before = Pt(40)

    # ==================== SLIDE 2: Executive Summary & Vision ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG_LIGHT)
    add_advanced_header(slide2, "Executive Summary & Core Strategic Vision")

    # 3 Grid Cards Layout
    card_width = Inches(3.64)
    card_height = Inches(4.8)
    card_top = Inches(1.9)

    cards_data_2 = [
        ("Platform Vision", "RetailPulse transforms raw retail data into predictive intelligence[cite: 3], bridging machine learning models with intuitive operational dashboards[cite: 2, 3] to minimize stock risk and enhance retention[cite: 2]."),
        ("Value Delivered", "Achieves high forecasting accuracy, mitigates inventory wastage, proactively identifies churn-prone accounts, and drives data-backed decision-making[cite: 3]."),
        ("Non-Functional Goals", "Engineered with a modular architecture[cite: 3], low-latency dashboard visualization[cite: 3], robust MLflow tracking[cite: 3], and automated containerized scalability[cite: 3].")
    ]

    for i, (head, body) in enumerate(cards_data_2):
        left_pos = Inches(0.8) + i * Inches(3.98)
        add_card(slide2, left_pos, card_top, card_width, card_height)
        
        # Inner text box
        tb = slide2.shapes.add_textbox(left_pos + Inches(0.3), card_top + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_h = tf.paragraphs[0]
        p_h.text = f"0{i+1}. {head}"
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_ACCENT_BLUE
        
        p_b = tf.add_paragraph()
        p_b.text = body
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(14)

    # ==================== SLIDE 3: End-to-End Architecture Workflow ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG_LIGHT)
    add_advanced_header(slide3, "System Architecture & MLOps Pipeline Flow")

    # 4 Structured Process Cards Layout (2x2 Grid)
    arch_steps = [
        ("1. Data Ingestion & Cleaning", "Processes multi-source transactional datasets, handles missing values, removes duplicates, and filters anomalies[cite: 3]."),
        ("2. Advanced Feature Engineering", "Extracts RFM metrics, time-series rolling statistics, customer attributes, and behavioral vectors[cite: 3]."),
        ("3. Predictive AI & Optimization", "Deploys K-Means segmentation, Churn classifiers, Prophet+LSTM hybrid forecasting, and EOQ inventory models[cite: 3]."),
        ("4. Deployment & MLOps Tracking", "Integrates multi-page Streamlit dashboards, Docker containerization, MLflow monitoring, and Airflow pipelines[cite: 3].")
    ]

    for i, (title, desc) in enumerate(arch_steps):
        col = i % 2
        row = i // 2
        left_pos = Inches(0.8) + col * Inches(5.9)
        top_pos = Inches(1.9) + row * Inches(2.4)
        
        add_card(slide3, left_pos, top_pos, Inches(5.6), Inches(2.1))
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.3), top_pos + Inches(0.25), Inches(5.0), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_before = Pt(8)

    # ==================== SLIDE 4: Execution Timeline & Milestones ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG_LIGHT)
    add_advanced_header(slide4, "Project Execution Timeline & Roadmap")

    weeks_data = [
        ("Week 1", "Planning & EDA", "Requirement analysis, dataset collection, cleaning pipelines, and Exploratory Data Analysis[cite: 3]."),
        ("Week 2", "Feature & Intelligence", "Feature engineering, K-Means RFM customer segmentation, and Churn prediction model implementation[cite: 3]."),
        ("Week 3", "Forecasting & Inventory", "Prophet & LSTM time-series modeling, hybrid integration, and inventory optimization rules[cite: 3]."),
        ("Week 4", "Deployment & Ops", "Streamlit dashboard development, Docker containerization, performance testing, and reporting[cite: 3].")
    ]

    card_w_4 = Inches(2.7)
    card_h_4 = Inches(4.8)
    for i, (wk, title, details) in enumerate(weeks_data):
        left_pos = Inches(0.8) + i * Inches(2.98)
        add_card(slide4, left_pos, Inches(1.9), card_w_4, card_h_4)
        
        tb = slide4.shapes.add_textbox(left_pos + Inches(0.25), Inches(2.1), card_w_4 - Inches(0.5), card_h_4 - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_wk = tf.paragraphs[0]
        p_wk.text = wk.upper()
        p_wk.font.name = 'Arial'
        p_wk.font.size = Pt(11)
        p_wk.font.bold = True
        p_wk.font.color.rgb = COLOR_ACCENT_CYAL
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_DARK
        p_t.space_before = Pt(6)
        
        p_d = tf.add_paragraph()
        p_d.text = details
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(12)

    # ==================== SLIDE 5: Technical Implementation & MLOps ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG_LIGHT)
    add_advanced_header(slide5, "Technical Highlights & MLOps Framework")

    sections_5 = [
        ("Advanced Modeling Techniques", "• Customer Segmentation via clustering and RFM behavioral scoring[cite: 3].\n• Churn Prediction using high-performance classification models[cite: 3].\n• Hybrid Demand Forecasting combining Prophet and LSTM architectures[cite: 3]."),
        ("Production MLOps Practices", "• MLflow for rigorous experiment parameter and metric tracking[cite: 3].\n• Optuna for automated hyperparameter tuning and optimization[cite: 3].\n• Data Drift Detection mechanisms and Airflow pipeline orchestration[cite: 3].")
    ]

    for i, (heading, content) in enumerate(sections_5):
        left_pos = Inches(0.8) + i * Inches(6.0)
        add_card(slide5, left_pos, Inches(1.9), Inches(5.733), Inches(4.8))
        
        tb = slide5.shapes.add_textbox(left_pos + Inches(0.4), Inches(2.2), Inches(4.933), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_h = tf.paragraphs[0]
        p_h.text = heading
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(18)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_ACCENT_BLUE
        
        p_c = tf.add_paragraph()
        p_c.text = content
        p_c.font.name = 'Arial'
        p_c.font.size = Pt(13)
        p_c.font.color.rgb = COLOR_TEXT_DARK
        p_c.space_before = Pt(16)

    # ==================== SLIDE 6: Platform Features Matrix (Table) ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG_LIGHT)
    add_advanced_header(slide6, "Enterprise Features & Acceptance Standards")

    rows, cols = 5, 3
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(4.733)

    headers = ["Module Feature", "Operational Description", "Acceptance Criteria Status"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BG_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.font.size = Pt(13)

    table_data = [
        ["Customer Segmentation", "Behavioral grouping via RFM & K-Means clustering[cite: 3]", "Customers successfully categorized[cite: 3]"],
        ["Demand & Hybrid Forecasting", "Predictive future trends using Prophet and LSTM models[cite: 3]", "Reliable high-accuracy forecast trends[cite: 3]"],
        ["Inventory Optimization", "Calculates safety stock levels and EOQ reorder points[cite: 3]", "Overstocking risk minimized[cite: 3]"],
        ["MLOps & Drift Detection", "Monitors experiment history and detects data distribution changes[cite: 3]", "Performance history & alerts active[cite: 3]"]
    ]

    for i, row_vals in enumerate(table_data):
        for j, val in enumerate(row_vals):
            cell = table.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_DARK

    # ==================== SLIDE 7: Technology Stack & Architecture Components ====================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_BG_LIGHT)
    add_advanced_header(slide7, "Technology Stack & Enterprise Architecture")

    tech_stack = [
        ("Core Language", "Python 3.12", "Primary environment for ML analytics and modeling[cite: 2, 3]."),
        ("Frontend UI", "Streamlit", "Interactive multi-page visual dashboards and control panels[cite: 2, 3]."),
        ("Machine Learning", "Scikit-Learn & TensorFlow", "Advanced classification, clustering, and neural networks[cite: 3]."),
        ("Containerization", "Docker & GitHub", "Portable execution environments and robust version control[cite: 2, 3].")
    ]

    card_w_7 = Inches(2.7)
    card_h_7 = Inches(4.8)
    for i, (cat, tech, desc) in enumerate(tech_stack):
        left_pos = Inches(0.8) + i * Inches(2.98)
        add_card(slide7, left_pos, Inches(1.9), card_w_7, card_h_7)
        
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.25), Inches(2.1), card_w_7 - Inches(0.5), card_h_7 - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_c = tf.paragraphs[0]
        p_c.text = cat.upper()
        p_c.font.name = 'Arial'
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_ACCENT_BLUE
        
        p_t = tf.add_paragraph()
        p_t.text = tech
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_DARK
        p_t.space_before = Pt(8)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(12)

    # ==================== SLIDE 8: Conclusion & Strategic Roadmap ====================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_BG_DARK)

    tb8 = slide8.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(5.0))
    tf8 = tb8.text_frame
    tf8.word_wrap = True

    p_head = tf8.paragraphs[0]
    p_head.text = "Conclusion & Future Roadmap"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(32)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEXT_LIGHT

    roadmap_points = [
        ("Practical Impact:", "Successfully demonstrated end-to-end data pipeline construction, predictive modeling, and professional application deployment[cite: 3]."),
        ("Industry Standards:", "Applied modular architecture, MLflow tracking, Docker containerization, and robust Git version control[cite: 3]."),
        ("Future Scale:", "Expansion toward real-time streaming analytics, cloud-native enterprise deployment, and automated recommendation engines[cite: 3].")
    ]

    for title, desc in roadmap_points:
        p = tf8.add_paragraph()
        p.space_before = Pt(20)
        r1 = p.add_run()
        r1.text = f"• {title} "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_ACCENT_CYAL
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_TEXT_LIGHT

    # Save Presentation File
    filename = "RetailPulse_Advanced_Presentation.pptx"
    prs.save(filename)
    print(f"Successfully generated advanced executive presentation: {filename}")

if __name__ == "__main__":
    create_advanced_retailpulse_ppt()